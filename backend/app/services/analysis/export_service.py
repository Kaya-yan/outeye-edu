"""
导出服务 - OutEye Edu 1.0

将教学方案导出为 PPT、Word 或 HTML 文档。
"""

import json
import os
import re
from datetime import datetime
from html import escape as html_escape
from typing import Dict, Any, List
from io import BytesIO
from loguru import logger


def export_pptx(plan: Dict[str, Any], title: str = "教学方案") -> BytesIO:
    """
    导出教学方案为 PPT

    Args:
        plan: 教学方案数据（含 difficulty_overview, teaching_suggestions, activity_designs, theoretical_basis）
        title: 演示文稿标题

    Returns:
        BytesIO 对象（PPT 文件）
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)  # Navy blue
    ACCENT = RGBColor(0xD4, 0x8B, 0x2C)   # Amber
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

    def add_title_slide():
        """封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PRIMARY
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        # Subtitle
        sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        gap = plan.get("learner_gap", {})
        p2.text = f"课文等级: {gap.get('text_level', '?')} | 学生水平: {gap.get('student_level', '?')} | 差距: {gap.get('gap', '?')}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p2.alignment = PP_ALIGN.CENTER

    def add_overview_slide():
        """课文难度概述"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_slide_header(slide, "课文难度概述")
        content = plan.get("difficulty_overview", "暂无概述")
        _add_body_text(slide, content)

    def add_suggestions_slide():
        """教学建议"""
        suggestions = plan.get("teaching_suggestions", [])
        if not suggestions:
            return
        # Split into slides of 4 items each
        for i in range(0, len(suggestions), 4):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_slide_header(slide, f"教学建议 ({i+1}-{min(i+4, len(suggestions))}/{len(suggestions)})")
            y_offset = Inches(1.8)
            for j, sug in enumerate(suggestions[i:i+4]):
                # 支持 dict 或 str
                if isinstance(sug, dict):
                    sug_text = sug.get("text", "")
                else:
                    sug_text = str(sug)
                # Number badge
                badge = slide.shapes.add_textbox(Inches(0.8), y_offset, Inches(0.5), Inches(0.4))
                tf = badge.text_frame
                p = tf.paragraphs[0]
                p.text = str(i + j + 1)
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT
                # Text
                txt = slide.shapes.add_textbox(Inches(1.5), y_offset, Inches(10.5), Inches(0.8))
                tf2 = txt.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = sug_text
                p2.font.size = Pt(16)
                p2.font.color.rgb = GRAY
                y_offset += Inches(1.2)

    def add_activities_slide():
        """课堂活动设计"""
        activities = plan.get("activity_designs", [])
        if not activities:
            return
        for i, act in enumerate(activities):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_slide_header(slide, f"活动 {i+1}: {act.get('name', '课堂活动')}")
            y_offset = Inches(1.8)
            for key, label in [("objective", "目标"), ("teacher_script", "教师指令"), ("steps", "步骤"), ("duration", "时间")]:
                val = act.get(key, "")
                if not val:
                    continue
                if isinstance(val, list):
                    val = "\n".join(str(v) for v in val)
                else:
                    val = str(val)
                # Label
                lbl = slide.shapes.add_textbox(Inches(0.8), y_offset, Inches(1.5), Inches(0.4))
                tf = lbl.text_frame
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = ACCENT
                # Value
                txt = slide.shapes.add_textbox(Inches(2.5), y_offset, Inches(9.5), Inches(0.8))
                tf2 = txt.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = val[:200]  # 截断超长文本
                p2.font.size = Pt(14)
                p2.font.color.rgb = GRAY
                y_offset += Inches(0.8)

    def add_theory_slide():
        """理论依据"""
        theory = plan.get("theoretical_basis", "")
        if not theory:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_slide_header(slide, "理论依据")
        _add_body_text(slide, theory)

    def _add_slide_header(slide, text: str):
        """添加幻灯片标题"""
        # Blue bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()
        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True

    def _add_body_text(slide, text: str):
        """添加正文内容"""
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)

    # Build slides
    add_title_slide()
    add_overview_slide()
    add_suggestions_slide()
    add_activities_slide()
    add_theory_slide()

    # Save to buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def export_docx(plan: Dict[str, Any], title: str = "教学方案") -> BytesIO:
    """
    导出教学方案为 Word 文档（结构化：概述/建议/活动含教师指令与rubric/分层/理论/词汇附录/文化附录）

    Args:
        plan: 教学方案数据（可附带 difficult_words、cultural_elements、differentiation）
        title: 文档标题

    Returns:
        BytesIO 对象（Word 文件）
    """
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_ALIGN_VERTICAL

    doc = Document()

    # 样式
    style = doc.styles["Normal"]
    style.font.size = Pt(11)
    style.font.name = "Microsoft YaHei"

    PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
    ACCENT = RGBColor(0xD4, 0x8B, 0x2C)
    MUTED = RGBColor(0x66, 0x66, 0x66)

    # 标题
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 副标题：等级 / 语言
    gap = plan.get("learner_gap", {})
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    gap_line = f"课文等级 {gap.get('text_level', '?')} · 学生水平 {gap.get('student_level', '?')} · 差距 {gap.get('gap', '?')}"
    if plan.get("language_name"):
        gap_line += f" · 语言 {plan['language_name']}"
    run = p.add_run(gap_line)
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED

    doc.add_paragraph()

    # 一、课文难度概述
    doc.add_heading("一、课文难度概述", level=1)
    overview = plan.get("difficulty_overview", "暂无概述")
    for line in str(overview).split("\n"):
        if line.strip():
            doc.add_paragraph(line.strip())

    # 二、教学建议
    suggestions = plan.get("teaching_suggestions", [])
    if suggestions:
        doc.add_heading("二、教学建议", level=1)
        for i, sug in enumerate(suggestions, 1):
            text = sug.get("text", "") if isinstance(sug, dict) else str(sug)
            hint = sug.get("data_hint", "") if isinstance(sug, dict) else ""
            p = doc.add_paragraph(style="List Number")
            p.add_run(text)
            if hint:
                run = p.add_run(f"  （{hint}）")
                run.font.size = Pt(9)
                run.font.color.rgb = MUTED

    # 三、课堂活动设计
    activities = plan.get("activity_designs", [])
    if activities:
        doc.add_heading("三、课堂活动设计", level=1)
        for i, act in enumerate(activities, 1):
            if not isinstance(act, dict):
                continue
            doc.add_heading(f"活动 {i}：{act.get('name', '课堂活动')}", level=2)

            # 元信息行：时长 + 形式
            meta_parts = []
            if act.get("duration"):
                meta_parts.append(f"时长 {act['duration']}")
            if act.get("interaction"):
                meta_parts.append(f"形式 {act['interaction']}")
            if meta_parts:
                p = doc.add_paragraph()
                run = p.add_run(" · ".join(meta_parts))
                run.italic = True
                run.font.size = Pt(10)
                run.font.color.rgb = MUTED

            # 目标
            obj = act.get("objective", "")
            if obj:
                p = doc.add_paragraph()
                r = p.add_run("目标：")
                r.bold = True
                r.font.color.rgb = ACCENT
                p.add_run(obj)

            # 教师指令原文（可直接念）
            teacher_script = act.get("teacher_script") or act.get("teacher_command") or ""
            if teacher_script:
                p = doc.add_paragraph()
                r = p.add_run("教师指令（可直接念）：")
                r.bold = True
                r.font.color.rgb = ACCENT
                p2 = doc.add_paragraph()
                p2.paragraph_format.left_indent = Inches(0.3)
                p2.add_run(teacher_script).italic = True

            # 步骤
            steps = act.get("steps", "")
            if steps:
                p = doc.add_paragraph()
                r = p.add_run("步骤：")
                r.bold = True
                r.font.color.rgb = ACCENT
                if isinstance(steps, str):
                    step_lines = [s.strip() for s in steps.split("\n") if s.strip()]
                else:
                    step_lines = [str(s) for s in steps if s]
                for j, step in enumerate(step_lines, 1):
                    sp = doc.add_paragraph(style="List Number")
                    sp.add_run(step)

            # 评估 rubric 表格
            rubric = act.get("rubric") or act.get("assessment_rubric")
            if rubric and isinstance(rubric, list):
                doc.add_paragraph().add_run("评估 Rubric：").bold = True
                # 表格：等级 / 描述 / 分值
                table = doc.add_table(rows=1 + len(rubric), cols=3)
                table.style = "Light Grid Accent 1"
                hdr = table.rows[0].cells
                hdr[0].text = "等级"
                hdr[1].text = "描述"
                hdr[2].text = "分值"
                for c in hdr:
                    for para in c.paragraphs:
                        for run in para.runs:
                            run.font.bold = True
                for r_idx, item in enumerate(rubric, 1):
                    if not isinstance(item, dict):
                        continue
                    row = table.rows[r_idx].cells
                    row[0].text = str(item.get("level", item.get("grade", "")))
                    row[1].text = str(item.get("description", item.get("desc", "")))
                    row[2].text = str(item.get("score", item.get("points", "")))

    # 四、分层教学（differentiation）
    diff = plan.get("differentiation", {})
    if diff and isinstance(diff, dict) and any(diff.values()):
        doc.add_heading("四、分层教学", level=1)
        for level_key, label in [("support", "基础层（需支持）"), ("core", "核心层（达标）"), ("challenge", "挑战层（拓展）")]:
            content = diff.get(level_key)
            if content:
                p = doc.add_paragraph()
                r = p.add_run(f"{label}：")
                r.bold = True
                r.font.color.rgb = ACCENT
                if isinstance(content, list):
                    for c in content:
                        doc.add_paragraph(str(c), style="List Bullet")
                else:
                    p.add_run(str(content))

    # 五、理论依据
    theory = plan.get("theoretical_basis", "")
    if theory:
        doc.add_heading("五、理论依据", level=1)
        for line in str(theory).split("\n"):
            if line.strip():
                doc.add_paragraph(line.strip())

    # 附录 A：词汇表
    difficult_words = plan.get("difficult_words", [])
    if difficult_words:
        doc.add_heading("附录 A：词汇表", level=1)
        table = doc.add_table(rows=1 + len(difficult_words), cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "词"
        hdr[1].text = "等级"
        hdr[2].text = "出现次数"
        hdr[3].text = "AWL"
        for c in hdr:
            for para in c.paragraphs:
                for run in para.runs:
                    run.font.bold = True
        for r_idx, w in enumerate(difficult_words, 1):
            if not isinstance(w, dict):
                continue
            row = table.rows[r_idx].cells
            row[0].text = str(w.get("word", ""))
            row[1].text = str(w.get("level", "—"))
            row[2].text = str(w.get("count", ""))
            row[3].text = "✓" if w.get("in_awl") else ""

    # 附录 B：文化元素
    cultural = plan.get("cultural_elements", [])
    if cultural:
        doc.add_heading("附录 B：文化元素", level=1)
        for i, e in enumerate(cultural, 1):
            if not isinstance(e, dict):
                continue
            p = doc.add_paragraph()
            r = p.add_run(f"{i}. [{e.get('category', '—')}] {e.get('keyword', '')}")
            r.bold = True
            r.font.color.rgb = ACCENT
            explanation = e.get("explanation", "")
            if explanation:
                doc.add_paragraph(explanation)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer


def export_html(plan: Dict[str, Any], title: str = "教学方案") -> BytesIO:
    """
    导出教学方案为交互式 HTML

    Args:
        plan: 教学方案数据（含 difficulty_overview, teaching_suggestions,
              activity_designs, theoretical_basis, learner_gap）
        title: 文档标题

    Returns:
        BytesIO 对象（HTML 文件）
    """
    # Load template
    template_path = os.path.join(
        os.path.dirname(__file__), "..", "..", "..", "templates", "html", "classroom_default.html"
    )
    template_path = os.path.normpath(template_path)

    try:
        with open(template_path, "r", encoding="utf-8") as f:
            template = f.read()
    except FileNotFoundError:
        logger.error(f"HTML 模板未找到: {template_path}")
        raise ValueError(f"HTML 模板文件不存在: {template_path}")

    # Build JSON data for injection
    gap = plan.get("learner_gap", {})
    json_data = _build_html_json(plan, title, gap)

    # Replace placeholders (HTML-escape title to prevent injection)
    html = template.replace("{{data.json}}", json.dumps(json_data, ensure_ascii=False, indent=2))
    html = html.replace("{{data.title}}", html_escape(title))

    # Write to buffer
    buffer = BytesIO(html.encode("utf-8"))
    return buffer


def _build_html_json(plan: Dict[str, Any], title: str, gap: Dict[str, Any]) -> Dict[str, Any]:
    """构建注入 HTML 模板的 JSON 数据"""

    # Parse suggestions
    suggestions = []
    for sug in plan.get("teaching_suggestions", []):
        if isinstance(sug, str):
            suggestions.append({"text": sug, "data_hint": ""})
        elif isinstance(sug, dict):
            suggestions.append({
                "text": sug.get("text", str(sug)),
                "data_hint": sug.get("data_hint", ""),
            })
        else:
            suggestions.append({"text": str(sug), "data_hint": ""})

    # Parse activities
    activities = []
    for act in plan.get("activity_designs", []):
        if not isinstance(act, dict):
            continue
        steps = act.get("steps", "")
        if isinstance(steps, str):
            steps = [s.strip() for s in steps.split("\n") if s.strip()]
        activities.append({
            "name": act.get("name", "课堂活动"),
            "icon": act.get("icon", "🎯"),
            "duration": act.get("duration", ""),
            "objective": act.get("objective", ""),
            "steps": steps,
            "data_hint": act.get("data_hint", ""),
        })

    # Parse theory
    theories = _parse_theory_text(plan.get("theoretical_basis", ""))

    return {
        "meta": {
            "title": title,
            "model": plan.get("model", "OutEye Edu"),
            "generated_at": datetime.now().strftime("%Y-%m-%d"),
            "level_from": gap.get("text_level", "B2"),
            "level_to": gap.get("student_level", "B1"),
            "gap": gap.get("gap", "i+1"),
            "tags": plan.get("tags", []),
        },
        "overview": {
            "summary": plan.get("difficulty_overview", ""),
            "flesch": plan.get("flesch_score", 0),
            "stats_hint": plan.get("stats_hint", ""),
        },
        "suggestions": suggestions,
        "activities": activities,
        "theories": theories,
        "data": {
            "vocabulary": plan.get("vocabulary", {}),
            "syntax": plan.get("syntax", {}),
            "discourse": plan.get("discourse", {}),
        },
    }


def _parse_theory_text(theory_text: str) -> List[Dict[str, Any]]:
    """从理论依据文本中解析出结构化数据"""
    if not theory_text:
        return []

    theories = []
    # Try to split by numbered patterns like "1." "2." or "一、" "二、"
    parts = re.split(r'\n(?=\d+[.、]|一、|二、|三、|四、)', theory_text)

    for part in parts:
        part = part.strip()
        if not part:
            continue
        # Extract name (first line) and description (rest)
        lines = part.split("\n", 1)
        name = re.sub(r'^[\d.、一二三四五六七八九十]+\s*', '', lines[0]).strip()
        desc = lines[1].strip() if len(lines) > 1 else ""

        theories.append({
            "name": name,
            "author": "",
            "description": desc or name,
            "tags": [],
        })

    # If no structured parts found, return as single entry
    if not theories:
        theories.append({
            "name": "理论依据",
            "author": "",
            "description": theory_text,
            "tags": [],
        })

    return theories
