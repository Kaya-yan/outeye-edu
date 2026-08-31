"""
课件 LLM 生成引擎（FIX-3 · F3.1/F3.2；HTML 链路 ④a 升级为三层架构）

HTML 链路（courseware_html_v2 · 三层）：框架层+主题层写死在
courseware_skeleton_v2.html（16:9 舞台、翻页/键盘/页码、交互行为、学术讲义
token 组），LLM 只生成逐页内容区 ```html 块 + 四选一强调色声明，后端拼装。
程序自检只查内容页（页数/单焦点/交互数/禁忌/行内色值），单页不合格定向
重生成 ≤2 轮，仍失败确定性净化；整副失败重试一次后回退 courseware_bootstrap，
fallback=True 由前端标注"简化版生成"，绝不静默。
PPT 链路（F3.3）：LLM 逐页大纲 JSON（≤6 要点/页、口语化讲者备注）
→ python-pptx 渲染 16:9；校验失败重试一次，仍失败回退确定性大纲。
Word 链路（F3.4）后续在此模块追加。
"""

from dataclasses import dataclass, field
from datetime import datetime
from html import escape as _html_escape
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple
from uuid import uuid4
from loguru import logger
import json
import re
import time

from app.services.prompt_manager import render_prompt, prompt_version
from app.services.analysis.fusion_generator import _esc, prepare_text
from app.services.courseware_themes import DEFAULT_THEME_ID, CoursewareTheme, get_theme

PROMPT_NAME = "courseware_html_v2"

_EXTERNAL_RE = re.compile(r"(?:src|href)\s*=\s*[\"']https?://|@import|<link[^>]+stylesheet", re.IGNORECASE)

# ============ ④a 三层架构：骨架（框架+主题）与内容层分离 ============

_SKELETON_PATH = Path(__file__).resolve().parent / "courseware_skeleton_v2.html"
_SKELETON_CACHE: Optional[str] = None

# 默认主题（academic）的基线 token；主题库见 courseware_themes.py
_ACADEMIC = get_theme(DEFAULT_THEME_ID)
THEME_PAPER = _ACADEMIC.tokens["paper"]
THEME_TOKENS = {k: _ACADEMIC.tokens[k] for k in ("paper", "ink", "text", "muted")}
ACCENT_PALETTE = {"#b5493e": "朱砂红", "#3e6b5a": "黛绿", "#99653a": "暖赭", "#35507a": "绀青"}
DEFAULT_ACCENT = _ACADEMIC.default_accent


@dataclass
class HTMLCoursewareResult:
    html: str
    editor_schema: Dict[str, Any]
    structure_sync: Dict[str, Any]
    self_check: Dict[str, Any] = field(default_factory=dict)
    prompt_version: str = ""
    model: str = ""
    fallback: bool = False
    retries: int = 0
    generation_duration: float = 0.0


def _format_plan_text(plan: Dict[str, Any]) -> str:
    """把结构化教案渲染为 LLM 可读的完整文本（课件内容的唯一来源）"""
    parts: List[str] = []

    if plan.get("framework"):
        parts += ["【教学设计框架】", str(plan["framework"])]
    objectives = plan.get("objectives") or []
    if objectives:
        parts.append("【教学目标】")
        for i, o in enumerate(objectives, 1):
            if isinstance(o, dict):
                line = f"目标{i}：{o.get('text', '')}"
                if o.get("bloom"):
                    line += f"（Bloom：{o['bloom']}）"
                if o.get("assessment"):
                    line += f"\n  评估方式：{o['assessment']}"
                parts.append(line)
            else:
                parts.append(f"目标{i}：{o}")
    if plan.get("difficulty_overview"):
        parts += ["【课文难度概述】", str(plan["difficulty_overview"])]

    suggestions = plan.get("teaching_suggestions") or []
    if suggestions:
        parts.append("【教学建议】")
        parts += [f"{i}. {s}" for i, s in enumerate(suggestions, 1)]

    parts.append("【课堂环节设计】")
    for i, act in enumerate(plan.get("activity_designs") or [], 1):
        if not isinstance(act, dict):
            parts.append(f"环节{i}：{act}")
            continue
        header = f"环节{i}：{act.get('name', '')}（{act.get('duration', '时长未标注')}）"
        parts.append(header)
        if act.get("objective"):
            parts.append(f"- 目标：{act['objective']}")
        if act.get("steps"):
            parts.append(f"- 步骤：{act['steps']}")
        if act.get("assessment"):
            parts.append(f"- 评估点：{act['assessment']}")

    assessment = plan.get("assessment") or {}
    formative = assessment.get("formative") or []
    summative = assessment.get("summative") or []
    if formative or summative:
        parts.append("【评估设计】")
        if formative:
            parts.append("形成性：" + "；".join(formative))
        if summative:
            parts.append("终结性：" + "；".join(summative))

    if plan.get("differentiation"):
        parts += ["【差异化教学策略】", str(plan["differentiation"])]
    if plan.get("theoretical_basis"):
        parts += ["【理论依据】", str(plan["theoretical_basis"])]

    return "\n".join(parts)


def _build_metrics_lines(analysis: Dict[str, Any]) -> str:
    """白盒关键指标精简行（供 LLM 取材，不重复教案已有内容）"""
    lines = []
    vocab = analysis.get("vocabulary") or {}
    difficult = ", ".join(
        d.get("word", "") for d in (vocab.get("difficult_words") or [])[:10]
    )
    if difficult:
        lines.append(f"- 难点词（前10）：{difficult}")
    if vocab.get("total_words"):
        lines.append(f"- 总词数 {vocab.get('total_words')}，不重复词 {vocab.get('unique_words')}")
    syntax = analysis.get("syntax") or {}
    max_sent = syntax.get("max_sentence") or {}
    if max_sent.get("preview"):
        lines.append(
            f"- 最长句（第{(max_sent.get('index') or 0) + 1}句，{max_sent.get('word_count')}词）：\"{str(max_sent.get('preview'))[:80]}\""
        )
    if syntax.get("avg_sentence_length"):
        lines.append(f"- 平均句长 {syntax['avg_sentence_length']} 词")
    discourse = analysis.get("discourse") or {}
    if discourse.get("genre_hint"):
        lines.append(f"- 体裁提示：{discourse['genre_hint']}，共 {discourse.get('paragraph_count', '?')} 段")
    return "\n".join(lines) or "-（无白盒指标）"


def _build_components_digest(components: List[Dict[str, Any]]) -> str:
    lines = []
    for c in components or []:
        level = c.get("interaction_level") or "static"
        lines.append(f"- {c.get('slug')} — {c.get('summary') or c.get('name')}（交互：{level}）")
    return "\n".join(lines) or "-（组件库为空，按约束规范自行设计原生交互）"


# ---- 内容层解析与程序自检（④a：框架项由骨架保证，只查内容页） ----

@dataclass
class _ContentPage:
    title: str
    intent: str
    html: str


_ACCENT_DECL_RE = re.compile(r"ACCENT\s*[:：]\s*(#[0-9a-fA-F]{6})")
_PAGE_META_RE = re.compile(r"<!--\s*page\s*[:：]\s*\d+\s*(?:\|\s*(.*?))?\s*-->", re.IGNORECASE)
_PAGE_INTENT_RE = re.compile(r"<!--\s*intent\s*[:：]\s*(.*?)\s*-->", re.IGNORECASE)
_RAW_COLOR_RE = re.compile(r"#[0-9a-fA-F]{3,8}\b|\brgb[a]?\s*\(|\bhsl[a]?\s*\(")
_GRADIENT_RE = re.compile(r"gradient\s*\(", re.IGNORECASE)
# ASCII 构造 emoji 区段，避免在源码嵌非 ASCII 字符区间（编辑易损坏）
_EMOJI_RE = re.compile(
    "["
    + chr(0x2600) + "-" + chr(0x27BF)
    + chr(0x2B00) + "-" + chr(0x2BFF)
    + chr(0xFE0F)
    + chr(0x1F000) + "-" + chr(0x1FAFF)
    + "]"
)
_FORBIDDEN_TAGS_RE = re.compile(r"<\s*(script|style|link|iframe|object|embed)\b", re.IGNORECASE)
_EVENT_ATTR_RE = re.compile(r"\son[a-z]+\s*=", re.IGNORECASE)
_INLINE_STYLE_FORBIDDEN_RE = re.compile(
    r"style\s*=\s*[\"'][^\"']*\b(color|background|font-family|line-height)", re.IGNORECASE
)
_PAGE_FOCUS_RE = re.compile(r"class\s*=\s*[\"'][^\"']*\bpage-focus\b")
_INTERACTION_MARKERS = {
    "reveal": re.compile(r"<details[^>]*class\s*=\s*[\"'][^\"']*\breveal\b", re.IGNORECASE),
    "timeline": re.compile(r"class\s*=\s*[\"'][^\"']*\btimeline\b", re.IGNORECASE),
    "vocab-card": re.compile(r"class\s*=\s*[\"'][^\"']*\bvocab-card\b", re.IGNORECASE),
    "timer": re.compile(r"class\s*=\s*[\"'][^\"']*\btimer\b|data-seconds\s*=", re.IGNORECASE),
}


def _parse_pages(answer: str) -> Tuple[str, List[_ContentPage]]:
    """解析内容层输出：ACCENT 声明 + 逐页 ```html 块。

    页注释（page/intent）契约在块内前两行；模型偏离写在围栏外时回看
    上一围栏结束到本围栏开始之间的文本兜底。
    """
    pages: List[_ContentPage] = []
    prev_end = 0
    for m_fence in re.finditer(r"```html\s*(.*?)\s*```", answer, re.DOTALL | re.IGNORECASE):
        block = m_fence.group(1)
        m_meta = _PAGE_META_RE.search(block[:400])
        m_intent = _PAGE_INTENT_RE.search(block[:400])
        content = block
        if m_meta or m_intent:
            for m in (m_meta, m_intent):
                if m:
                    content = content.replace(m.group(0), "", 1)
        else:
            lookback = answer[prev_end:m_fence.start()][-400:]
            m_meta = _PAGE_META_RE.search(lookback)
            m_intent = _PAGE_INTENT_RE.search(lookback)
        pages.append(
            _ContentPage(
                title=(m_meta.group(1).strip() if m_meta and m_meta.group(1) else ""),
                intent=(m_intent.group(1).strip() if m_intent else ""),
                html=content.strip(),
            )
        )
        prev_end = m_fence.end()
    m_accent = _ACCENT_DECL_RE.search(answer)
    accent = (m_accent.group(1).lower() if m_accent else "")
    return accent, pages


def _interaction_types(pages: List[_ContentPage]) -> Set[str]:
    all_html = "\n".join(p.html for p in pages)
    return {name for name, rx in _INTERACTION_MARKERS.items() if rx.search(all_html)}


def _validate_content_page(page: _ContentPage) -> List[str]:
    """单页程序自检：单焦点 / 无脚本 / 无事件属性 / 无行内色值 / 无渐变 / 无 emoji / 无外链"""
    problems = []
    focus_count = len(_PAGE_FOCUS_RE.findall(page.html))
    if focus_count != 1:
        problems.append(f"必须恰好一个 .page-focus，实际 {focus_count} 个")
    if len(page.html) < 120:
        problems.append("内容过短（<120 字符），疑似截断")
    if _FORBIDDEN_TAGS_RE.search(page.html):
        problems.append("含禁用标签（script/style/link/iframe 等），行为与样式由骨架负责")
    if _EVENT_ATTR_RE.search(page.html):
        problems.append("含事件属性（on*=）")
    if _RAW_COLOR_RE.search(page.html):
        problems.append("含行内色值（#hex/rgb()/hsl()），颜色只能用 var(--token)")
    if _GRADIENT_RE.search(page.html):
        problems.append("含渐变")
    if _EMOJI_RE.search(page.html):
        problems.append("含 emoji 或装饰性符号")
    if _INLINE_STYLE_FORBIDDEN_RE.search(page.html):
        problems.append("行内 style 含禁用属性（color/background/font-family/line-height）")
    if _EXTERNAL_RE.search(page.html):
        problems.append("含外链资源，违反单文件约束")
    return problems


def _validate_deck(pages: List[_ContentPage], min_pages: int) -> Optional[str]:
    """整副课件结构性校验：返回 None 通过，否则失败原因（触发整体重试/回退）"""
    if len(pages) < min_pages:
        return f"页面数不足：需 ≥{min_pages} 页，实际 {len(pages)}"
    types = _interaction_types(pages)
    if len(types) < 3:
        return f"交互类型不足：需 ≥3 种（reveal/timeline/vocab-card/timer），实际 {sorted(types) or '无'}"
    return None


def _sanitize_page(html_str: str) -> str:
    """重生成仍失败时的兜底净化：确定性剥除脚本/事件/违禁行内样式/emoji，保证硬性红线"""
    s = re.sub(r"<script\b.*?</script>|<style\b.*?</style>|<iframe\b.*?</iframe>", "", html_str, flags=re.IGNORECASE | re.DOTALL)
    s = re.sub(r"<\s*(link|object|embed)\b[^>]*/?>", "", s, flags=re.IGNORECASE)
    s = re.sub(r"\son[a-z]+\s*=\s*(\"[^\"]*\"|'[^']*'|[^\s>]+)", "", s, flags=re.IGNORECASE)
    s = re.sub(r"style\s*=\s*(\"[^\"]*\b(?:color|background|font-family|line-height|gradient)[^\"]*\"|'[^']*\b(?:color|background|font-family|line-height|gradient)[^']*')", "", s, flags=re.IGNORECASE)
    s = _EMOJI_RE.sub("", s)
    return s.strip()


def _load_skeleton() -> str:
    global _SKELETON_CACHE
    if _SKELETON_CACHE is None:
        _SKELETON_CACHE = _SKELETON_PATH.read_text(encoding="utf-8")
    return _SKELETON_CACHE


def _blend_colors(hex_fg: str, hex_bg: str, ratio: float) -> str:
    fg = [int(hex_fg[i:i + 2], 16) for i in (1, 3, 5)]
    bg = [int(hex_bg[i:i + 2], 16) for i in (1, 3, 5)]
    mixed = [round(ratio * fg[i] + (1 - ratio) * bg[i]) for i in range(3)]
    return "#" + "".join(f"{c:02x}" for c in mixed)


def _assemble_skeleton(title: str, accent: str, pages: List[_ContentPage], theme: Optional[CoursewareTheme] = None) -> str:
    """三层拼装：骨架（框架）+ 主题 token 组与微调（④b）+ 逐页内容 section（重编页码）"""
    theme = theme or _ACADEMIC
    sections = []
    for i, p in enumerate(pages, 1):
        attrs = f'<section class="page" data-page="{i}" data-title="{_html_escape(p.title or f"第{i}页", quote=True)}"'
        if p.intent:
            attrs += f' data-intent="{_html_escape(p.intent, quote=True)}"'
        sections.append(f'{attrs}>\n{p.html}\n</section>')
    token_block = (
        theme.token_css()
        + f"\n  --accent:{accent};\n  --accent-soft:{_blend_colors(accent, theme.tokens['paper'], 0.12)};"
    )
    html = _load_skeleton().replace("<!--SECTIONS-->", "\n".join(sections))
    return (
        html.replace("__TITLE__", _html_escape(title))
        .replace("__TOKEN_BLOCK__", token_block)
        .replace("__THEME_EXTRA__", theme.extra_css)
    )


def _relative_luminance(hex_color: str) -> float:
    def channel(c: int) -> float:
        v = c / 255
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4

    r, g, b = (int(hex_color[i:i + 2], 16) for i in (1, 3, 5))
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def contrast_ratio(hex_a: str, hex_b: str) -> float:
    """WCAG 对比度；主题 token 组合在测试中断言 ≥4.5:1（内容层禁行内色值即承袭该保证）"""
    la, lb = _relative_luminance(hex_a), _relative_luminance(hex_b)
    lighter, darker = max(la, lb), min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


def _regen_page(generator: Any, system_prompt: str, user_prompt: str, page: _ContentPage, index: int, problems: List[str]) -> Optional[_ContentPage]:
    """单页定向重生成：成功且通过自检才替换，否则返回 None（调用方保留原稿）"""
    user = (
        f"{user_prompt}\n\n"
        f"以下是课件第 {index} 页（标题：{page.title or '无'}，教学意图：{page.intent or '无'}）的现有内容：\n"
        f"```html\n{page.html}\n```\n\n"
        "程序自检发现以下问题：\n" + "\n".join(f"- {p}" for p in problems) + "\n\n"
        "请只重写这一页的内容区 HTML，修复上述全部问题，教学设计保持不变。"
        "输出：一个 ```html 代码块（块内前两行仍是页注释 page/intent），此外不输出任何文字。"
        "提醒：颜色只用 var(--ink/--text/--muted/--paper/--accent/--line)；禁止 script/style/link/iframe、"
        "事件属性、行内色值、渐变、emoji、外链；每页恰好一个 .page-focus。"
    )
    answer, _usage = generator._generate_with_api(
        [{"role": "system", "content": system_prompt}, {"role": "user", "content": user}]
    )
    _, pages = _parse_pages(answer)
    if len(pages) != 1:
        return None
    if _validate_content_page(pages[0]):
        return None
    return pages[0]


def _build_prompt(
    *,
    title: str,
    plan: Dict[str, Any],
    analysis: Dict[str, Any],
    text: str,
    language_name: str,
    text_level: str,
    student_level: str,
    duration_minutes: int,
    course_type: str,
    class_size: int,
    native_language: str,
    components: List[Dict[str, Any]],
    theme: Optional[CoursewareTheme] = None,
) -> str:
    stages = plan.get("activity_designs") or []
    cw_theme = theme or _ACADEMIC
    _, user_prompt = render_prompt(
        PROMPT_NAME,
        title=_esc(title),
        language_name=_esc(language_name),
        text_level=_esc(text_level),
        student_level=_esc(student_level),
        duration_minutes=int(duration_minutes or 90),
        course_type=_esc(course_type or "综合"),
        class_size=int(class_size or 30),
        native_language=_esc(native_language or "中文"),
        full_text=_esc(prepare_text(text or "")),
        plan_text=_esc(_format_plan_text(plan)),
        metrics_lines=_esc(_build_metrics_lines(analysis)),
        components_digest=_esc(_build_components_digest(components)),
        theme_desc=f"「{cw_theme.name}」主题（{cw_theme.palette_desc}）",
    )
    return user_prompt


def _wrap_llm_schema(title: str, html: str, source_meta: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """LLM 产物的编辑器 schema：单页 html_embed，编辑器渲染时收编 rendered_html"""
    page_id = f"page-{uuid4().hex[:8]}"
    return {
        "meta": {
            "title": title,
            "mode": "slides",
            "template_id": "classroom_default",
            "source_type": "from_plan_llm",
            "created_at": datetime.utcnow().isoformat(),
            "source_meta": source_meta or {},
        },
        "outline": [{"id": page_id, "label": "AI 生成课件", "source_key": "rendered_html"}],
        "pages": [
            {
                "id": page_id,
                "kind": "slide",
                "title": title,
                "blocks": [
                    {
                        "id": f"block-{uuid4().hex[:8]}",
                        "type": "html_embed",
                        "label": "AI 生成内容",
                        "editable": "free",
                        "source_key": "rendered_html",
                        "content": {"html": html},
                    }
                ],
            }
        ],
    }


def _structure_sync_from_pages(schema: Dict[str, Any]) -> Dict[str, Any]:
    pages = []
    for page in schema.get("pages", []):
        pages.append(
            {
                "page_id": page.get("id"),
                "page_title": page.get("title"),
                "blocks": [
                    {
                        "block_id": b.get("id"),
                        "source_key": b.get("source_key"),
                        "editable": b.get("editable"),
                        "type": b.get("type"),
                    }
                    for b in page.get("blocks", [])
                ],
            }
        )
    return {"pages": pages}


def generate_html_courseware(
    *,
    title: str,
    plan: Dict[str, Any],
    analysis: Dict[str, Any],
    text: str,
    language_name: str = "英语",
    text_level: str = "",
    student_level: str = "",
    duration_minutes: int = 90,
    course_type: Optional[str] = None,
    class_size: Optional[int] = None,
    native_language: Optional[str] = None,
    components: Optional[List[Dict[str, Any]]] = None,
    learner_gap: Optional[Dict[str, Any]] = None,
    enhancement_tags: Optional[List[str]] = None,
    theme: Optional[str] = None,
) -> HTMLCoursewareResult:
    """
    生成单文件交互 HTML 课件（④a 三层架构）。

    内容层 LLM 输出（逐页 ```html 块 + ACCENT 声明）经程序自检：
    整副失败自动重试一次（携带原因）；单页不合格定向重生成（≤2 轮），
    仍失败则确定性净化保底线；全部失败回退 courseware_bootstrap，fallback=True。
    """
    start_time = time.time()
    version = prompt_version(PROMPT_NAME)
    components = components or []
    cw_theme = get_theme(theme)
    min_pages = max(3, len(plan.get("activity_designs") or []) + 1)

    model_name = "template-fallback"
    fallback_used = True
    retries = 0
    html = ""
    accent = DEFAULT_ACCENT
    pages: List[_ContentPage] = []
    regenerated: List[int] = []
    sanitized: List[int] = []
    self_check: Dict[str, Any] = {}
    raw_answer = ""

    try:
        user_prompt = _build_prompt(
            title=title,
            plan=plan,
            analysis=analysis or {},
            text=text,
            language_name=language_name,
            text_level=text_level,
            student_level=student_level,
            duration_minutes=duration_minutes,
            course_type=course_type,
            class_size=class_size,
            native_language=native_language,
            components=components,
            theme=cw_theme,
        )

        from app.services.rag import RAGGenerator
        from app.core.config import settings

        model_name = getattr(settings, "LLM_MODEL", "deepseek-chat")
        generator = RAGGenerator(
            api_key=getattr(settings, "LLM_API_KEY", None),
            api_base=getattr(settings, "LLM_BASE_URL", None),
            model=model_name,
            max_tokens=8000,
            temperature=0.7,
        )

        if generator.use_api:
            system_prompt, _ = render_prompt(PROMPT_NAME)
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ]

            answer, _usage = generator._generate_with_api(messages)
            raw_answer = answer
            accent, pages = _parse_pages(answer)
            reason = _validate_deck(pages, min_pages)

            if reason:
                # F3.5：整副失败自动重试一次，携带失败原因
                retries = 1
                logger.warning(f"HTML 课件内容层首次校验失败（{reason}），重试一次")
                messages += [
                    {"role": "assistant", "content": answer[-2000:]},
                    {"role": "user", "content": f"上一次输出未通过结构校验：{reason}。请重新按输出契约输出：ACCENT 首行声明 + 逐页 ```html 块 + 自检 JSON。"},
                ]
                answer, _usage = generator._generate_with_api(messages)
                raw_answer = answer
                accent, pages = _parse_pages(answer)
                reason = _validate_deck(pages, min_pages)

            if reason:
                logger.warning(f"HTML 课件重试仍失败（{reason}），回退模板拼装")
            else:
                fallback_used = False

                # 逐页程序自检：定向重生成 ≤2 轮，仍失败则确定性净化保硬性红线
                for round_no in range(2):
                    failing = [(i, probs) for i, p in enumerate(pages) if (probs := _validate_content_page(p))]
                    if not failing:
                        break
                    if len(failing) > 4:
                        reason = "不合格页过多（>4），整体质量不足：" + "；".join(
                            f"第{i + 1}页({'/'.join(probs[:2])})" for i, probs in failing[:3]
                        )
                        logger.warning(f"HTML 课件{reason}，回退模板拼装")
                        fallback_used = True
                        break
                    logger.info(f"内容页自检第{round_no + 1}轮：第 {[i + 1 for i, _ in failing]} 页需重写")
                    for i, probs in failing:
                        fixed = _regen_page(generator, system_prompt, user_prompt, pages[i], i + 1, probs)
                        if fixed is not None:
                            pages[i] = fixed
                            if i + 1 not in regenerated:
                                regenerated.append(i + 1)

                if not fallback_used:
                    residual = [(i, probs) for i, p in enumerate(pages) if (probs := _validate_content_page(p))]
                    for i, probs in residual:
                        logger.warning(f"第{i + 1}页重生成后仍不合格（{'；'.join(probs)}），确定性净化保底线")
                        pages[i].html = _sanitize_page(pages[i].html)
                        sanitized.append(i + 1)

                if not fallback_used:
                    accent_note = None
                    if accent not in ACCENT_PALETTE:
                        logger.warning(f"强调色 {accent or '未声明'} 不在色板内，回退默认 {DEFAULT_ACCENT}")
                        accent_note = f"声明值 {accent or '未声明'} 不在色板，已用默认 {DEFAULT_ACCENT}"
                        accent = DEFAULT_ACCENT
                    self_check = {
                        "prompt_version": version,
                        "accent": accent,
                        "accent_note": accent_note,
                        "theme": cw_theme.id,
                        "pages_count": len(pages),
                        "page_intents": [
                            {"page": i + 1, "title": p.title, "intent": p.intent} for i, p in enumerate(pages)
                        ],
                        "interaction_types": sorted(_interaction_types(pages)),
                        "regenerated_pages": regenerated,
                        "sanitized_pages": sanitized,
                        "llm_self_check": _extract_selfcheck(raw_answer),
                    }
        else:
            logger.warning("LLM 不可用，HTML 课件回退模板拼装")
    except Exception as e:
        logger.warning(f"HTML 课件 LLM 生成失败，回退模板拼装: {e}")

    if fallback_used:
        from app.services.courseware_bootstrap import build_courseware_from_plan

        bootstrap = build_courseware_from_plan(
            title=title,
            mode="slides",
            template_id="classroom_default",
            plan=plan,
            learner_gap=learner_gap,
            enhancement_tags=enhancement_tags,
            components=components,
        )
        html = bootstrap["rendered_html"]
        schema = bootstrap["editor_schema_json"]
        sync = bootstrap["structure_sync_json"]
        self_check = {
            "prompt_version": "fallback",
            "notes": "LLM 生成不可用或未通过校验，已用教案模板拼装（简化版生成）",
        }
        source_meta = {"generated_by": "template_fallback", "prompt_version": version}
    else:
        html = _assemble_skeleton(title, accent, pages, theme=cw_theme)
        source_meta = {"generated_by": "llm_html_v2", "prompt_version": version, "theme": cw_theme.id}
        schema = _wrap_llm_schema(title, html, source_meta)
        sync = _structure_sync_from_pages(schema)

    return HTMLCoursewareResult(
        html=html,
        editor_schema=schema,
        structure_sync=sync,
        self_check=self_check,
        prompt_version=version,
        model=model_name if not fallback_used else "template-fallback",
        fallback=fallback_used,
        retries=retries,
        generation_duration=round(time.time() - start_time, 2),
    )


def _extract_selfcheck(answer: str) -> Dict[str, Any]:
    from app.services.analysis.fusion_generator import _extract_self_check

    return _extract_self_check(answer)


# ============ F3.3 PPT 链路 ============

PPT_PROMPT_NAME = "courseware_ppt_v1"

_PPT_KINDS = {"cover", "agenda", "content", "quote", "interaction", "vocab", "summary"}

# Morandi 平台色（与前端 tokens 同源）
_PPT_INK = "3A3A37"
_PPT_INK_SOFT = "6B6B66"
_PPT_SAGE = "96A790"
_PPT_SAGE_DARK = "767870"
_PPT_ACCENT = "D8C46A"
_PPT_CANVAS = "FAF9F6"
_PPT_ROSE_LIGHT = "EFE7E2"


@dataclass
class PPTCoursewareResult:
    outline: Dict[str, Any]
    pptx_bytes: BytesIO
    slide_count: int
    self_check: Dict[str, Any] = field(default_factory=dict)
    prompt_version: str = ""
    model: str = ""
    fallback: bool = False
    retries: int = 0
    generation_duration: float = 0.0


def _extract_json_object(answer: str) -> Dict[str, Any]:
    """提取 ```json 代码块；无围栏时退而取首个平衡的 {...} 片段"""
    blocks = re.findall(r"```json\s*(.*?)\s*```", answer, re.DOTALL | re.IGNORECASE)
    for block in blocks:
        try:
            return json.loads(block)
        except json.JSONDecodeError:
            continue
    depth = 0
    start = -1
    for i, ch in enumerate(answer):
        if ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start >= 0:
                try:
                    return json.loads(answer[start : i + 1])
                except json.JSONDecodeError:
                    start = -1
    return {}


def _validate_ppt_outline(outline: Dict[str, Any], min_slides: int) -> Optional[str]:
    """返回 None 表示通过，否则返回失败原因（用于重试反馈与回退判定）"""
    slides = outline.get("slides")
    if not isinstance(slides, list) or len(slides) < min_slides:
        actual = len(slides) if isinstance(slides, list) else 0
        return f"页数不足：需 ≥{min_slides} 页（环节数+3），实际 {actual}"
    problems = []
    for i, s in enumerate(slides, 1):
        if not isinstance(s, dict):
            problems.append(f"第{i}页不是对象")
            continue
        if not str(s.get("title", "")).strip():
            problems.append(f"第{i}页缺标题")
        bullets = s.get("bullets")
        if not isinstance(bullets, list) or not bullets:
            problems.append(f"第{i}页无要点")
        elif len(bullets) > 6:
            problems.append(f"第{i}页要点 {len(bullets)} 条，超过 6 条上限")
        if len(str(s.get("notes", "")).strip()) < 20:
            problems.append(f"第{i}页讲者备注过短")
    if problems:
        return "；".join(problems[:5])
    return None


def _build_ppt_prompt(
    *,
    title: str,
    plan: Dict[str, Any],
    analysis: Dict[str, Any],
    text: str,
    language_name: str,
    text_level: str,
    student_level: str,
    duration_minutes: int,
    course_type: str,
    class_size: int,
    native_language: str,
    slide_count_hint: int,
) -> str:
    _, user_prompt = render_prompt(
        PPT_PROMPT_NAME,
        title=_esc(title),
        language_name=_esc(language_name),
        text_level=_esc(text_level),
        student_level=_esc(student_level),
        duration_minutes=int(duration_minutes or 90),
        course_type=_esc(course_type or "综合"),
        class_size=int(class_size or 30),
        native_language=_esc(native_language or "中文"),
        slide_count_hint=int(slide_count_hint),
        full_text=_esc(prepare_text(text or "")),
        plan_text=_esc(_format_plan_text(plan)),
        metrics_lines=_esc(_build_metrics_lines(analysis)),
    )
    return user_prompt


def _fallback_ppt_outline(title: str, plan: Dict[str, Any], language_name: str, duration_minutes: int) -> Dict[str, Any]:
    """LLM 不可用时的确定性大纲：封面/目标/逐环节/总结，notes 标注模板生成"""
    slides: List[Dict[str, Any]] = [
        {
            "kind": "cover",
            "title": title,
            "bullets": [f"{language_name} · 课堂课件", f"课时 {duration_minutes or 90} 分钟"],
            "notes": "封面页。开场问候后快速过页，报出本课主题与课时安排。（模板生成：AI 大纲暂不可用，可生成后人工调整）",
            "layout_hint": "center",
        }
    ]
    objectives = plan.get("objectives") or []
    if objectives:
        slides.append({
            "kind": "agenda",
            "title": "教学目标",
            "bullets": [str(o.get("text", o) if isinstance(o, dict) else o)[:40] for o in objectives[:6]],
            "notes": "目标页。逐条口头展开，强调本课结束时学生能做到什么。（模板生成）",
            "layout_hint": "bullets",
        })
    for i, act in enumerate(plan.get("activity_designs") or [], 1):
        if not isinstance(act, dict):
            continue
        bullets: List[str] = []
        if act.get("objective"):
            bullets.append(str(act["objective"])[:40])
        if act.get("steps"):
            steps = re.split(r"[；;。]", str(act["steps"]))
            bullets += [s.strip()[:30] for s in steps if s.strip()][:4]
        slides.append({
            "kind": "content",
            "title": f"{act.get('name', f'环节{i}')[:15]}",
            "bullets": bullets[:6] or ["（本环节要点待补充）"],
            "notes": f"本环节约 {act.get('duration', '—')}。{act.get('assessment') or act.get('objective') or ''}（模板生成）"[:200],
            "layout_hint": "bullets",
        })
    slides.append({
        "kind": "summary",
        "title": "总结与作业",
        "bullets": [str(s)[:30] for s in (plan.get("assessment", {}) or {}).get("summative", [])][:5] or ["回顾本课要点"],
        "notes": "总结页。回收本课目标达成情况并布置作业。（模板生成）",
        "layout_hint": "bullets",
    })
    return {"slides": slides, "self_check": {"source": "template_fallback"}}


def _render_pptx(outline: Dict[str, Any], title: str) -> BytesIO:
    """把大纲 JSON 渲染为 16:9 .pptx（Morandi 版式，讲者备注入 notes）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(hexstr: str) -> "RGBColor":
        return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

    ink, ink_soft = rgb(_PPT_INK), rgb(_PPT_INK_SOFT)
    sage, sage_dark, accent = rgb(_PPT_SAGE), rgb(_PPT_SAGE_DARK), rgb(_PPT_ACCENT)
    canvas, rose_light = rgb(_PPT_CANVAS), rgb(_PPT_ROSE_LIGHT)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def fill_bg(slide, color) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def add_rect(slide, x, y, w, h, color) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()

    def add_text(slide, x, y, w, h, lines, size, color, *, bold=False, italic=False, align=PP_ALIGN.LEFT, space_after=10):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(line)
            p.alignment = align
            p.space_after = Pt(space_after)
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.color.rgb = color
                run.font.bold = bold
                run.font.italic = italic
        return box

    slides = outline.get("slides", [])
    total = max(len(slides), 1)
    for idx, s in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        kind = s.get("kind") if s.get("kind") in _PPT_KINDS else "content"
        bullets = [str(b) for b in (s.get("bullets") or [])][:6]
        hint = s.get("layout_hint")
        title_text = str(s.get("title", "")).strip() or "—"
        center = kind in {"cover", "quote", "interaction"} or hint == "center"

        if kind == "interaction":
            fill_bg(slide, rose_light)
            add_text(slide, 1.0, 1.6, 11.3, 1.2, [title_text], 30, ink, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, 1.5, 3.1, 10.3, 3.4, bullets, 20, ink_soft, align=PP_ALIGN.CENTER, space_after=14)
        elif kind == "quote":
            fill_bg(slide, canvas)
            add_rect(slide, 0, 0, 13.333, 0.12, sage)
            add_text(slide, 1.0, 0.8, 11.3, 0.8, [title_text], 14, sage_dark, align=PP_ALIGN.CENTER)
            add_text(slide, 1.8, 2.4, 9.7, 3.6, ["“ " + b + " ”" for b in bullets] or ["…"], 22, ink_soft, italic=True, align=PP_ALIGN.CENTER, space_after=16)
        elif kind == "cover":
            fill_bg(slide, canvas)
            add_rect(slide, 0, 6.9, 13.333, 0.6, sage)
            add_text(slide, 1.0, 2.4, 11.3, 1.6, [title_text], 36, ink, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, 1.0, 4.2, 11.3, 1.6, bullets, 16, ink_soft, align=PP_ALIGN.CENTER, space_after=8)
        else:
            fill_bg(slide, canvas)
            add_rect(slide, 0, 0, 13.333, 0.12, sage)
            add_text(slide, 0.8, 0.5, 11.7, 1.0, [title_text], 24, ink, bold=True)
            add_rect(slide, 0.85, 1.45, 1.2, 0.06, accent)
            if hint == "two_col" and len(bullets) >= 4:
                half = (len(bullets) + 1) // 2
                add_text(slide, 0.8, 2.0, 5.7, 4.6, bullets[:half], 18, ink_soft, space_after=12)
                add_text(slide, 6.9, 2.0, 5.7, 4.6, bullets[half:], 18, ink_soft, space_after=12)
            else:
                add_text(slide, 0.8, 2.0, 11.7, 4.6, bullets, 18, ink_soft, space_after=12)

        notes = str(s.get("notes", "")).strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        add_text(slide, 0.8, 7.05, 8.0, 0.4, [title[:40]], 9, ink_soft)
        add_text(slide, 11.5, 7.05, 1.4, 0.4, [f"{idx} / {total}"], 9, ink_soft, align=PP_ALIGN.RIGHT)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_ppt_courseware(
    *,
    title: str,
    plan: Dict[str, Any],
    analysis: Dict[str, Any],
    text: str,
    language_name: str = "英语",
    text_level: str = "",
    student_level: str = "",
    duration_minutes: int = 90,
    course_type: Optional[str] = None,
    class_size: Optional[int] = None,
    native_language: Optional[str] = None,
) -> PPTCoursewareResult:
    """
    生成 16:9 课堂放映 PPT：LLM 逐页大纲 JSON → python-pptx 渲染。
    大纲校验失败自动重试一次（携带原因）；仍失败回退确定性大纲，fallback=True。
    """
    start_time = time.time()
    version = prompt_version(PPT_PROMPT_NAME)
    activities = plan.get("activity_designs") or []
    min_slides = max(5, len(activities) + 3)
    slide_count_hint = len(activities) + 4

    model_name = "template-fallback"
    fallback_used = True
    retries = 0
    outline: Dict[str, Any] = {}
    self_check: Dict[str, Any] = {}
    raw_answer = ""

    try:
        user_prompt = _build_ppt_prompt(
            title=title,
            plan=plan,
            analysis=analysis or {},
            text=text,
            language_name=language_name,
            text_level=text_level,
            student_level=student_level,
            duration_minutes=duration_minutes,
            course_type=course_type or "综合",
            class_size=class_size or 30,
            native_language=native_language or "中文",
            slide_count_hint=slide_count_hint,
        )

        from app.services.rag import RAGGenerator
        from app.core.config import settings

        model_name = getattr(settings, "LLM_MODEL", "deepseek-chat")
        generator = RAGGenerator(
            api_key=getattr(settings, "LLM_API_KEY", None),
            api_base=getattr(settings, "LLM_BASE_URL", None),
            model=model_name,
            max_tokens=6000,
            temperature=0.6,
        )

        if generator.use_api:
            system_prompt, _ = render_prompt(PPT_PROMPT_NAME)
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ]

            answer, _usage = generator._generate_with_api(messages)
            raw_answer = answer
            outline = _extract_json_object(answer)
            reason = _validate_ppt_outline(outline, min_slides)

            if reason:
                retries = 1
                logger.warning(f"PPT 大纲首次校验失败（{reason}），重试一次")
                messages += [
                    {"role": "assistant", "content": answer[-2000:]},
                    {"role": "user", "content": f"上一次输出未通过结构校验：{reason}。请重新输出完整的大纲 JSON，严格遵循输出契约。"},
                ]
                answer, _usage = generator._generate_with_api(messages)
                raw_answer = answer
                outline = _extract_json_object(answer)
                reason = _validate_ppt_outline(outline, min_slides)

            if reason:
                logger.warning(f"PPT 大纲重试仍失败（{reason}），回退确定性大纲")
            else:
                fallback_used = False
                self_check = outline.get("self_check") or {}
        else:
            logger.warning("LLM 不可用，PPT 回退确定性大纲")
    except Exception as e:
        logger.warning(f"PPT 课件 LLM 生成失败，回退确定性大纲: {e}")

    if fallback_used:
        outline = _fallback_ppt_outline(title, plan, language_name, duration_minutes)
        self_check = {"prompt_version": "fallback", "notes": "LLM 大纲不可用或未通过校验，已用教案确定性大纲渲染（简化版生成）"}

    pptx_bytes = _render_pptx(outline, title)
    return PPTCoursewareResult(
        outline=outline,
        pptx_bytes=pptx_bytes,
        slide_count=len(outline.get("slides", [])),
        self_check=self_check,
        prompt_version=version,
        model=model_name if not fallback_used else "template-fallback",
        fallback=fallback_used,
        retries=retries,
        generation_duration=round(time.time() - start_time, 2),
    )


# ============ F3.4 Word 链路 ============

WORD_PROMPT_NAME = "courseware_word_v1"

_WORD_KINDS = {"cover", "objectives", "stage", "board", "homework", "appendix"}


@dataclass
class WordCoursewareResult:
    outline: Dict[str, Any]
    docx_bytes: BytesIO
    section_count: int
    self_check: Dict[str, Any] = field(default_factory=dict)
    prompt_version: str = ""
    model: str = ""
    fallback: bool = False
    retries: int = 0
    generation_duration: float = 0.0


def _validate_word_outline(outline: Dict[str, Any], min_sections: int) -> Optional[str]:
    """返回 None 表示通过，否则返回失败原因（用于重试反馈与回退判定）"""
    sections = outline.get("sections")
    if not isinstance(sections, list) or len(sections) < min_sections:
        actual = len(sections) if isinstance(sections, list) else 0
        return f"章节数不足：需 ≥{min_sections} 节（环节数+3），实际 {actual}"
    problems = []
    for i, s in enumerate(sections, 1):
        if not isinstance(s, dict):
            problems.append(f"第{i}节不是对象")
            continue
        if not str(s.get("heading", "")).strip():
            problems.append(f"第{i}节缺标题")
        has_content = bool(s.get("bullets") or s.get("paragraphs") or s.get("table"))
        if not has_content:
            problems.append(f"第{i}节无内容（bullets/paragraphs/table 至少一项）")
    if problems:
        return "；".join(problems[:5])
    return None


def _build_word_prompt(
    *,
    title: str,
    plan: Dict[str, Any],
    analysis: Dict[str, Any],
    text: str,
    language_name: str,
    text_level: str,
    student_level: str,
    duration_minutes: int,
    course_type: str,
    class_size: int,
    native_language: str,
) -> str:
    _, user_prompt = render_prompt(
        WORD_PROMPT_NAME,
        title=_esc(title),
        language_name=_esc(language_name),
        text_level=_esc(text_level),
        student_level=_esc(student_level),
        duration_minutes=int(duration_minutes or 90),
        course_type=_esc(course_type or "综合"),
        class_size=int(class_size or 30),
        native_language=_esc(native_language or "中文"),
        full_text=_esc(prepare_text(text or "")),
        plan_text=_esc(_format_plan_text(plan)),
        metrics_lines=_esc(_build_metrics_lines(analysis)),
    )
    return user_prompt


def _fallback_word_outline(title: str, plan: Dict[str, Any], language_name: str, duration_minutes: int) -> Dict[str, Any]:
    """LLM 不可用时的确定性文档结构：封面/目标/逐环节/板书/作业"""
    sections: List[Dict[str, Any]] = [
        {
            "kind": "cover",
            "heading": title,
            "bullets": [f"语种：{language_name}", f"课时：{duration_minutes or 90} 分钟"],
            "paragraphs": ["（模板生成：AI 文档结构暂不可用，以下为教案确定性文档，可人工调整。）"],
        },
        {
            "kind": "objectives",
            "heading": "教学目标",
            "bullets": [str(o.get("text", o) if isinstance(o, dict) else o) for o in (plan.get("objectives") or [])][:8],
            "paragraphs": [],
        },
    ]
    board_lines = []
    for i, act in enumerate(plan.get("activity_designs") or [], 1):
        if not isinstance(act, dict):
            continue
        steps = re.split(r"[；;]", str(act.get("steps", "")))
        step_lines = [f"步骤{j}：{s.strip()}" for j, s in enumerate(steps, 1) if s.strip()]
        sections.append({
            "kind": "stage",
            "heading": f"环节{i} {act.get('name', '')}（{act.get('duration', '—')}）",
            "bullets": ([f"目标：{act['objective']}"] if act.get("objective") else [])
                       + ([f"评估点：{act['assessment']}"] if act.get("assessment") else []),
            "paragraphs": step_lines or ["（步骤待补充）"],
        })
        board_lines.append(f"环节{i} {act.get('name', '')}：{act.get('objective', '')}")
    sections.append({
        "kind": "board",
        "heading": "板书设计",
        "bullets": board_lines or ["（板书待设计）"],
        "paragraphs": [],
    })
    sections.append({
        "kind": "homework",
        "heading": "作业与课后评估",
        "bullets": [str(s) for s in (plan.get("assessment", {}) or {}).get("summative", [])] or ["（作业待布置）"],
        "paragraphs": [],
    })
    return {"sections": sections, "self_check": {"source": "template_fallback"}}


def _render_docx(outline: Dict[str, Any], title: str) -> BytesIO:
    """把文档结构 JSON 渲染为 .docx（封面 + 逐节标题/条目/段落/表格）"""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, 0)

    for s in outline.get("sections", []):
        kind = s.get("kind") if s.get("kind") in _WORD_KINDS else "stage"
        heading = str(s.get("heading", "")).strip()
        if kind != "cover" and heading:
            doc.add_heading(heading, level=1)
        for para in s.get("paragraphs") or []:
            doc.add_paragraph(str(para))
        for bullet in s.get("bullets") or []:
            doc.add_paragraph(str(bullet), style="List Bullet")
        table = s.get("table")
        if isinstance(table, dict) and isinstance(table.get("headers"), list):
            headers = [str(h) for h in table["headers"]]
            rows = [[str(c) for c in row] for row in (table.get("rows") or [])]
            t = doc.add_table(rows=1, cols=len(headers))
            t.style = "Table Grid"
            for j, h in enumerate(headers):
                cell = t.rows[0].cells[j]
                cell.text = h
                for run in cell.paragraphs[0].runs:
                    run.font.bold = True
            for row in rows:
                cells = t.add_row().cells
                for j, val in enumerate(row[: len(headers)]):
                    cells[j].text = val

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def generate_word_courseware(
    *,
    title: str,
    plan: Dict[str, Any],
    analysis: Dict[str, Any],
    text: str,
    language_name: str = "英语",
    text_level: str = "",
    student_level: str = "",
    duration_minutes: int = 90,
    course_type: Optional[str] = None,
    class_size: Optional[int] = None,
    native_language: Optional[str] = None,
) -> WordCoursewareResult:
    """
    生成教师课堂执行文档（Word）：LLM 结构 JSON → python-docx 渲染。
    校验失败自动重试一次；仍失败回退确定性结构，fallback=True。
    """
    start_time = time.time()
    version = prompt_version(WORD_PROMPT_NAME)
    activities = plan.get("activity_designs") or []
    min_sections = len(activities) + 3

    model_name = "template-fallback"
    fallback_used = True
    retries = 0
    outline: Dict[str, Any] = {}
    self_check: Dict[str, Any] = {}
    raw_answer = ""

    try:
        user_prompt = _build_word_prompt(
            title=title,
            plan=plan,
            analysis=analysis or {},
            text=text,
            language_name=language_name,
            text_level=text_level,
            student_level=student_level,
            duration_minutes=duration_minutes,
            course_type=course_type or "综合",
            class_size=class_size or 30,
            native_language=native_language or "中文",
        )

        from app.services.rag import RAGGenerator
        from app.core.config import settings

        model_name = getattr(settings, "LLM_MODEL", "deepseek-chat")
        generator = RAGGenerator(
            api_key=getattr(settings, "LLM_API_KEY", None),
            api_base=getattr(settings, "LLM_BASE_URL", None),
            model=model_name,
            max_tokens=6000,
            temperature=0.5,
        )

        if generator.use_api:
            system_prompt, _ = render_prompt(WORD_PROMPT_NAME)
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ]

            answer, _usage = generator._generate_with_api(messages)
            raw_answer = answer
            outline = _extract_json_object(answer)
            reason = _validate_word_outline(outline, min_sections)

            if reason:
                retries = 1
                logger.warning(f"Word 结构首次校验失败（{reason}），重试一次")
                messages += [
                    {"role": "assistant", "content": answer[-2000:]},
                    {"role": "user", "content": f"上一次输出未通过结构校验：{reason}。请重新输出完整的文档结构 JSON，严格遵循输出契约。"},
                ]
                answer, _usage = generator._generate_with_api(messages)
                raw_answer = answer
                outline = _extract_json_object(answer)
                reason = _validate_word_outline(outline, min_sections)

            if reason:
                logger.warning(f"Word 结构重试仍失败（{reason}），回退确定性结构")
            else:
                fallback_used = False
                self_check = outline.get("self_check") or {}
        else:
            logger.warning("LLM 不可用，Word 回退确定性结构")
    except Exception as e:
        logger.warning(f"Word 课件 LLM 生成失败，回退确定性结构: {e}")

    if fallback_used:
        outline = _fallback_word_outline(title, plan, language_name, duration_minutes)
        self_check = {"prompt_version": "fallback", "notes": "LLM 结构不可用或未通过校验，已用教案确定性文档渲染（简化版生成）"}

    docx_bytes = _render_docx(outline, title)
    return WordCoursewareResult(
        outline=outline,
        docx_bytes=docx_bytes,
        section_count=len(outline.get("sections", [])),
        self_check=self_check,
        prompt_version=version,
        model=model_name if not fallback_used else "template-fallback",
        fallback=fallback_used,
        retries=retries,
        generation_duration=round(time.time() - start_time, 2),
    )
